"""
Что содержит: функции извлечения текста из PPTX, PDF и markdown-представления слайдов.
За что отвечает: за чтение исходных файлов презентации и получение текстового содержимого по слайдам.
Где используется: импортируется в `src.app.services.processor` как первый этап обработки презентации.
"""

from __future__ import annotations

import re
from pathlib import Path


def load_pptx_slides(pptx_path: str | Path) -> list[str]:
    """Извлекает текст из каждого слайда PPTX-файла."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError("Для чтения PPTX-файлов требуется пакет `python-pptx`.") from exc

    presentation = Presentation(str(pptx_path))
    slides: list[str] = []
    for slide in presentation.slides:
        text_fragments: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    text_fragments.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            text_fragments.append(cell_text)
        slides.append("\n".join(text_fragments).strip())
    return slides


def load_markdown_slides(markdown_path: str | Path) -> list[str]:
    """Читает markdown-файл со слайдами и возвращает их по порядку."""
    markdown_text = Path(markdown_path).read_text(encoding="utf-8")
    pattern = re.compile(
        r"^###\s*Слайд\s*(\d+)\s*\n(.*?)(?=^###\s*Слайд\s*\d+\s*$|\Z)",
        re.IGNORECASE | re.MULTILINE | re.DOTALL,
    )
    slides: list[tuple[int, str]] = []
    for match in pattern.finditer(markdown_text):
        slide_number = int(match.group(1))
        slide_body = match.group(2).strip()
        slides.append((slide_number, slide_body))
    if not slides:
        raise ValueError("Markdown-файл должен содержать слайды в формате `### Слайд N`.")
    return [slide_text for _, slide_text in sorted(slides, key=lambda slide: slide[0])]


def load_pdf_slides(pdf_path: str | Path) -> list[str]:
    """Извлекает текст из страниц PDF и возвращает его постранично."""
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ImportError("Для чтения PDF-файлов требуется пакет `pypdf`.") from exc

    reader = PdfReader(str(pdf_path))
    slides = [(page.extract_text() or "").strip() for page in reader.pages]
    if not slides:
        raise ValueError("PDF-файл не содержит страниц для обработки.")
    return slides
